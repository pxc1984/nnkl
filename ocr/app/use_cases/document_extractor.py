"""Lightweight native text extraction for office documents."""

from __future__ import annotations

import re
from pathlib import Path

import fitz
from docx import Document as DocxDocument
from pptx import Presentation

_ALLOWED_CONTROL_CHARS = {"\n", "\r", "\t"}
_CONTROL_CHAR_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]")


class UnsupportedDocumentTypeError(Exception):
    """Raised when the document type is not supported."""


def extract_native_document_text(file_path: Path, *, output_format: str) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        text = _extract_pdf_text(file_path)
    elif suffix == ".docx":
        text = _extract_docx_text(file_path)
    elif suffix == ".pptx":
        text = _extract_pptx_text(file_path)
    else:
        raise UnsupportedDocumentTypeError(f"Unsupported document type: {suffix or 'unknown'}")

    if output_format == "latex":
        return _to_latex(text)
    return _to_markdown(text)


def has_native_pdf_text(file_path: Path, *, min_chars: int = 50) -> bool:
    """Heuristic: enough extractable text on the first page means a native-text PDF."""
    try:
        with fitz.open(file_path) as document:
            if document.page_count == 0:
                return False
            return len(sanitize_extracted_text(document[0].get_text()).strip()) >= min_chars
    except Exception:  # noqa: BLE001
        return False


def should_use_native_pdf_text(file_path: Path, *, min_chars: int = 50) -> bool:
    """Uses native extraction only when the text layer looks usable."""
    try:
        with fitz.open(file_path) as document:
            if document.page_count == 0:
                return False
            raw_text = document[0].get_text()
            sanitized = sanitize_extracted_text(raw_text).strip()
            return len(sanitized) >= min_chars and _is_reasonable_pdf_text(raw_text, sanitized)
    except Exception:  # noqa: BLE001
        return False


def _extract_pdf_text(file_path: Path) -> str:
    pages: list[str] = []
    with fitz.open(file_path) as document:
        for page in document:
            text = sanitize_extracted_text(page.get_text()).strip()
            if text:
                pages.append(text)
    return "\n\n".join(pages).strip()


def sanitize_extracted_text(text: str) -> str:
    """Removes NUL and other disallowed control characters from extracted text."""
    if not text:
        return ""
    return "".join(char for char in text if char >= " " or char in _ALLOWED_CONTROL_CHARS)


def _is_reasonable_pdf_text(raw_text: str, sanitized_text: str) -> bool:
    if not sanitized_text:
        return False

    control_chars = len(_CONTROL_CHAR_RE.findall(raw_text))
    if control_chars / max(len(raw_text), 1) > 0.01:
        return False

    printable_chars = sum(1 for char in sanitized_text if not char.isspace())
    if printable_chars == 0:
        return False

    alpha_chars = sum(1 for char in sanitized_text if char.isalpha())
    return (alpha_chars / printable_chars) >= 0.2


def _extract_docx_text(file_path: Path) -> str:
    document = DocxDocument(file_path)
    lines = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                lines.append(" | ".join(cells))

    return "\n\n".join(lines).strip()


def _extract_pptx_text(file_path: Path) -> str:
    presentation = Presentation(file_path)
    slides: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = [f"Slide {index}"]
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            text = text.strip()
            if text:
                parts.append(text)
        if len(parts) > 1:
            slides.append("\n\n".join(parts))

    return "\n\n".join(slides).strip()


def _to_markdown(text: str) -> str:
    return text


def _to_latex(text: str) -> str:
    replacements = {
        "\\": r"\textbackslash{}",
        "&": r"\&",
        "%": r"\%",
        "$": r"\$",
        "#": r"\#",
        "_": r"\_",
        "{": r"\{",
        "}": r"\}",
        "~": r"\textasciitilde{}",
        "^": r"\textasciicircum{}",
    }
    escaped = "".join(replacements.get(char, char) for char in text)
    paragraphs = [part.strip() for part in escaped.split("\n\n") if part.strip()]
    return "\n\n".join(paragraphs)
