"""Pytest fixtures."""

from __future__ import annotations

import os
import sys
from collections.abc import Generator
from io import BytesIO
from pathlib import Path
from unittest.mock import MagicMock

import pytest
from fastapi.testclient import TestClient

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

# Лёгкие моки mineru CLI для unit-тестов без тяжёлых зависимостей
if "mineru" not in sys.modules:
    sys.modules.setdefault("mineru", MagicMock())

TEST_ROOT = Path(__file__).resolve().parent / "fixtures"
TEST_DATABASE = TEST_ROOT / "test.sqlite3"

os.environ.setdefault("DATABASE_URL", f"sqlite:///{TEST_DATABASE}")
os.environ.setdefault("TEMP_DIR", str(TEST_ROOT / "tmp"))
os.environ.setdefault("MINERU_USE_GPU", "false")


@pytest.fixture(autouse=True)
def setup_dirs() -> Generator[None, None, None]:
    TEST_ROOT.mkdir(parents=True, exist_ok=True)
    if TEST_DATABASE.exists():
        TEST_DATABASE.unlink()
    yield
    if TEST_DATABASE.exists():
        TEST_DATABASE.unlink()


@pytest.fixture
def client() -> Generator[TestClient, None, None]:
    from app.main import app

    with TestClient(app) as test_client:
        yield test_client


@pytest.fixture
def sample_pdf() -> Path:
    """Минимальный валидный PDF для тестов."""
    pdf_path = TEST_ROOT / "sample.pdf"
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text(
        (72, 72),
        "Тестовый материал: сталь 09Г2С, предел прочности σ = 450 МПа при 20°C. "
        "Марка стали соответствует ГОСТ 19281.",
    )
    doc.save(pdf_path)
    doc.close()
    return pdf_path


@pytest.fixture
def sample_docx_bytes() -> bytes:
    from docx import Document

    document = Document()
    document.add_paragraph("DOCX title")
    document.add_paragraph("DOCX body text")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Cell A"
    table.rows[0].cells[1].text = "Cell B"

    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


@pytest.fixture
def sample_pptx_bytes() -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "PPTX title"
    slide.placeholders[1].text = "PPTX bullet"

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


@pytest.fixture
def db_session(client: TestClient):
    session = client.app.state.session_factory()
    try:
        yield session
    finally:
        session.close()
